from pptx import Presentation


def create_presentation(slides):

    prs = Presentation()

    for slide in slides:

        layout = prs.slide_layouts[1]

        slide_obj = prs.slides.add_slide(layout)

        title = slide_obj.shapes.title
        content = slide_obj.placeholders[1]

        title.text = slide["title"]
        content.text = "\n".join(slide["points"])

    file_name = "presentation.pptx"

    prs.save(file_name)

    return file_name
