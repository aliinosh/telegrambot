from pptx.util import Pt


def build_slide(prs, title, points):

    layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for point in points:

        paragraph = text_frame.add_paragraph()
        paragraph.text = point
        paragraph.level = 0
        paragraph.font.size = Pt(24)

    return slide
