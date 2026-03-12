from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def add_chart(slide):

    chart_data = ChartData()

    chart_data.categories = ["A", "B", "C"]

    chart_data.add_series(
        "Example",
        (10, 20, 30)
    )

    slide.shapes.add_chart(

        XL_CHART_TYPE.COLUMN_CLUSTERED,

        Inches(1),
        Inches(2),
        Inches(6),
        Inches(4),

        chart_data
    )
