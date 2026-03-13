from pptx import Presentation
from config.constants import APP_NAME, MAX_SLIDES


def create_presentation(title="Presentation"):
    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = APP_NAME

    return prs
